import re
import os
from openai import OpenAI
from pptx import Presentation
from urllib.parse import quote
from app.config import API_KEY, MODEL_EMBEDDING, API_BASE_URL

client = OpenAI(api_key=API_KEY.strip())

# ==========================================
# 1. 임베딩 및 텍스트 유틸리티
# ==========================================
def get_embedding(text):
    text = text.replace("\n", " ")
    if not text.strip(): return []
    try:
        return client.embeddings.create(input=[text], model=MODEL_EMBEDDING).data[0].embedding
    except: return []

def chunk_text(text, chunk_size=1000, overlap=100):
    """
    일반 텍스트(PPT 등)를 위한 단순 청킹 함수
    """
    if not text: return []
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        if end == text_len: break
        start += (chunk_size - overlap)
    return chunks

def extract_text_from_pptx(pptx_path):
    """
    PPTX 파일에서 텍스트 추출
    """
    try:
        prs = Presentation(pptx_path)
        full_text = []
        for i, slide in enumerate(prs.slides):
            text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            full_text.append(f"\n[Slide {i+1}] {text}")
        return "\n".join(full_text)
    except: return ""

# ==========================================
# 2. Markdown 파싱 로직
# ==========================================
def parse_markdown_structure(text):
    header_match = re.search(r'#\s+항목\s+정보:\s+\[([^\]]+)\]', text)
    global_key = header_match.group(1).strip() if header_match else "unknown"
    
    parts = global_key.split('_')
    result = {
        "process_id": parts[0].strip() if len(parts) > 0 else "unknown",
        "process_name": parts[1].strip() if len(parts) > 1 else "unknown",
        "global_key": global_key,
        "sections": {"기본정보": [], "상세속성": [], "연관항목": []},
        "chunks": []
    }

    section_blocks = re.split(r'##\s+\[([^\]]+)\]\s+항목:\s+\[([^\]]+)\]', text)
    
    for i in range(1, len(section_blocks), 3):
        sec_name = section_blocks[i].strip()
        sec_key = section_blocks[i+1].strip()
        content = section_blocks[i+2].strip()

        result["chunks"].append({
            "text": f"[{sec_key}] [{sec_name}]\n{content}",
            "section": sec_name,
            "global_key": sec_key
        })

        if "상세속성" in sec_name:
            attr_blocks = re.split(r'####\s+', content)
            for ab in attr_blocks:
                lines = ab.strip().split('\n')
                if len(lines) > 1:
                    result["sections"]["상세속성"].append({"key": lines[0].strip(), "value": "\n".join(lines[1:]).strip()})

        elif "연관항목" in sec_name:
            blocks = re.split(r'-\s+연관\s+No\.', content)
            for b in blocks:
                if not b.strip(): continue
                id_match = re.search(r'Identity\)\*\*:\s*\[([^\]]+)\]', b)
                if id_match:
                    g_key = id_match.group(1).strip()
                    name_parts = g_key.split('_')
                    rel_name = name_parts[1] if len(name_parts) > 1 else g_key
                    
                    dynamic_rels = re.findall(r'-\s+\*\*([^*]+)\*\*:\s*([^\n]+)', b)
                    rel_label = "연관"
                    exclude = ["시스템 ID", "System ID", "항목 식별 정보", "Item Identity", "바로가기"]
                    for k, v in dynamic_rels:
                        if not any(ex in k for ex in exclude):
                            rel_label = k.strip(); break

                    result["sections"]["연관항목"].append({
                        "name": rel_name,
                        "global_key": g_key,
                        "rel_label": rel_label
                    })
    return result